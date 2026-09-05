import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide

    # Color Palette (Dark High-Tech Physics Lab Theme)
    BG_COLOR = RGBColor(11, 19, 43)        # Deep Navy
    CARD_BG = RGBColor(28, 37, 65)         # Slate Navy
    ACCENT_CYAN = RGBColor(72, 202, 228)   # Electric Blue
    ACCENT_PURPLE = RGBColor(114, 9, 183)  # Violet
    TEXT_WHITE = RGBColor(248, 249, 250)   # Clean White
    TEXT_MUTED = RGBColor(148, 163, 184)   # Slate Light
    ACCENT_GREEN = RGBColor(6, 214, 160)   # Emerald Green
    ACCENT_AMBER = RGBColor(255, 183, 3)   # Warm Amber

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background() # No border
        return bg

    def add_header(slide, title_text, category_text="WLCG CRIC TECHNICAL POC"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = ACCENT_CYAN

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.7))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, title, items, badge_color=ACCENT_CYAN, badge_text=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = RGBColor(45, 55, 85)
        shape.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.5), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = badge_color
        p0.space_after = Pt(12)

        for itm in items:
            p = tf.add_paragraph()
            p.text = f"• {itm}"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(8)

    # ----------------- SLIDE 1: Title Slide -----------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.2))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    
    p_tag = tf1.paragraphs[0]
    p_tag.text = "CERN COMPUTING RESOURCE INFORMATION CATALOGUE (CRIC) POC"
    p_tag.font.size = Pt(13)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_CYAN
    p_tag.space_after = Pt(10)

    p_main = tf1.add_paragraph()
    p_main.text = "WLCG Infrastructure Explorer"
    p_main.font.size = Pt(40)
    p_main.font.bold = True
    p_main.font.color.rgb = TEXT_WHITE
    p_main.space_after = Pt(12)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Explainable Multi-Source Grid Reconciliation & Provenance Architecture"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = ACCENT_CYAN
    p_sub.space_after = Pt(24)

    # Meta card at bottom
    meta_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.8))
    meta_box.fill.solid()
    meta_box.fill.fore_color.rgb = CARD_BG
    meta_box.line.color.rgb = ACCENT_CYAN
    meta_box.line.width = Pt(1)

    m_tf = meta_box.text_frame
    m_tf.word_wrap = True
    mp1 = m_tf.paragraphs[0]
    mp1.text = "Technical Architecture Demonstration | IT-CE-LCG-2026-54-GRAP"
    mp1.font.size = Pt(14)
    mp1.font.bold = True
    mp1.font.color.rgb = TEXT_WHITE
    mp1.space_after = Pt(8)

    mp2 = m_tf.add_paragraph()
    mp2.text = "Key Pillars: 1. Raw Provenance Preservation  |  2. Explainable Evidence Scoring  |  3. Data Quality Awareness"
    mp2.font.size = Pt(12)
    mp2.font.color.rgb = ACCENT_GREEN
    mp2.space_after = Pt(8)

    mp3 = m_tf.add_paragraph()
    mp3.text = "Data Sources: GOCDB (REST/XML) • BDII (LDAP/GLUE2) • OSG Topology (XML/JSON)"
    mp3.font.size = Pt(11)
    mp3.font.color.rgb = TEXT_MUTED

    # ----------------- SLIDE 2: Problem Statement -----------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "The Challenge: Heterogeneous Grid Resource Catalogues")

    add_card(s2, 0.8, 1.8, 3.6, 4.8, "1. Disparate Protocols", [
        "GOCDB serves authoritative European site registries via REST XML.",
        "BDII relies on LDAP queries exposing complex GLUE2 hierarchical schemas.",
        "OSG Topology exposes US grid resource groups via distinct XML/JSON endpoints.",
        "No single protocol or unified key connects entities out-of-the-box."
    ], ACCENT_CYAN)

    add_card(s2, 4.8, 1.8, 3.6, 4.8, "2. Schema Drift & Flaws", [
        "Whitespace inconsistencies (e.g. leading spaces in GOCDB GIIS URLs).",
        "Conflicting naming conventions across domains (FQDN vs local alias).",
        "Legacy records with missing geographic coordinates or stale endpoints.",
        "Traditional ETL pipelines fail silently or drop non-compliant records."
    ], ACCENT_AMBER)

    add_card(s2, 8.8, 1.8, 3.6, 4.8, "3. Black-Box Entity Merging", [
        "Heuristics without audit trails create confusion for Grid operators.",
        "Unexplainable ML deduplication obscures root-cause discrepancy sources.",
        "Critical need for transparent, explainable scoring criteria.",
        "WLCG operations demand reproducible provenance down to raw timestamps."
    ], ACCENT_PURPLE)

    # ----------------- SLIDE 3: Architectural Overview -----------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "Architecture: 4-Stage Resilient Processing Pipeline")

    add_card(s3, 0.8, 1.8, 2.7, 4.8, "Stage 1: Ingestion", [
        "Independent Adapters for GOCDB, BDII, and OSG.",
        "Immutable Raw Snapshot capture.",
        "No destructive transformation at capture point.",
        "Timestamped provenance hash generation."
    ], ACCENT_CYAN)

    add_card(s3, 3.8, 1.8, 2.7, 4.8, "Stage 2: Sanitization", [
        "Defensive field normalization.",
        "Whitespace trimming and host canonicalization.",
        "Zero data loss philosophy.",
        "Records flagged as valid_with_warnings rather than dropped."
    ], ACCENT_GREEN)

    add_card(s3, 6.8, 1.8, 2.7, 4.8, "Stage 3: Reconciliation", [
        "Weighted evidence scoring engine.",
        "+40 Domain ID match.",
        "+25 Endpoint Host match.",
        "+20 ISO Country match.",
        "Deterministic confidence classification."
    ], ACCENT_AMBER)

    add_card(s3, 9.8, 1.8, 2.7, 4.8, "Stage 4: Unified Model", [
        "Canonical WLCG Site topology.",
        "Multi-provider side-by-side view.",
        "Interactive Explorer dashboard.",
        "Netlify Edge distribution with live SSR/API endpoints."
    ], ACCENT_PURPLE)

    # ----------------- SLIDE 4: Raw Provenance -----------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Architectural Pillar 1: Preserving Raw Provenance")

    add_card(s4, 0.8, 1.8, 5.6, 4.8, "Non-Destructive Ingestion", [
        "Raw payloads from GOCDB, BDII LDAP, and OSG are preserved byte-for-byte.",
        "No downstream reconciliation step overwrites or alters original records.",
        "Every reconciled field retains a pointer back to its specific source record ID.",
        "Guarantees that auditability and legal data provenance remain 100% intact.",
        "Permits retroactive re-execution of new matching rules without re-fetching historical snapshots."
    ], ACCENT_CYAN)

    add_card(s4, 6.8, 1.8, 5.6, 4.8, "Provenance Tree Model", [
        "Interactive field-level drill-down for grid operations engineers.",
        "Answers 'Where did this endpoint URL originate?' in 1 click.",
        "Example Provenance Path:",
        "  Canonical Site -> Endpoint",
        "  ↳ Source: BDII GLUE2 (sbdii.ui.savba.sk:2170)",
        "  ↳ Entity: GLUE2EndpointID=iisas-ce01.savba.sk:8443",
        "  ↳ Ingestion Timestamp: 2026-09-05T08:30:00Z",
        "  ↳ Raw Snapshot Hash: sha256:7f49a2b8e3c1..."
    ], ACCENT_GREEN)

    # ----------------- SLIDE 5: Data Quality & Sanitization -----------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Architectural Pillar 2: Data Quality Awareness")

    add_card(s5, 0.8, 1.8, 5.6, 4.8, "Handling Upstream Real-World Anomalies", [
        "Real-world grid catalogues frequently contain edge-case anomalies.",
        "Observed GOCDB Case: GIIS URL fields contain leading/trailing whitespaces.",
        "Observed BDII Case: Mismatched LDAP DN casing and legacy schema artifacts.",
        "Observed OSG Case: Resource Group alias collisions with legacy FQDNs.",
        "Our Approach: Automatic defensive sanitation without discarding records."
    ], ACCENT_AMBER)

    add_card(s5, 6.8, 1.8, 5.6, 4.8, "The 'valid_with_warnings' Contract", [
        "Traditional pipelines treat warnings as fatal errors or ignore them silently.",
        "CRIC POC introduces explicit validity classification:",
        "  • VALID: Clean record conforming to canonical schema.",
        "  • VALID_WITH_WARNINGS: Sanitized anomaly preserved with warning badge.",
        "  • INVALID: Critical structural failure with logged diagnostic.",
        "Exposes quality dashboards so site managers can correct source metadata."
    ], ACCENT_CYAN)

    # ----------------- SLIDE 6: Explainable Entity Resolution -----------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Architectural Pillar 3: Explainable Entity Resolution")

    add_card(s6, 0.8, 1.8, 5.6, 4.8, "Weighted Evidence Scoring Engine", [
        "Reconciliation requires auditable mathematical criteria, not black-box guesses.",
        "Evidence Breakdown:",
        "  • +40 points: Canonical Domain / FQDN identifier match",
        "  • +25 points: Compute/Storage Endpoint Hostname match",
        "  • +20 points: ISO 3166 Country Code match",
        "Confidence Tier Thresholds:",
        "  • High Confidence: Score >= 60 (Automatic Merge)",
        "  • Review Required: 30 <= Score < 60 (Operator Attention)",
        "  • Distinct: Score < 30 (Separate Site Entities)"
    ], ACCENT_PURPLE)

    add_card(s6, 6.8, 1.8, 5.6, 4.8, "Example Resolution: IISAS-Bratislava", [
        "GOCDB Record: Site ID 41 | IISAS-Bratislava | Slovakia",
        "BDII Record: GLUE2DomainID=IISAS-Bratislava | Slovakia",
        "Score Calculation:",
        "  [✓] Domain ID Match ('IISAS-Bratislava'): +40",
        "  [✓] Endpoint Host Match ('ce01.savba.sk'): +25",
        "  [✓] Country Match ('Slovakia'): +20",
        "Total Resolution Score = 85 (HIGH CONFIDENCE)",
        "Full transparent breakdown visible to grid operators in the UI."
    ], ACCENT_GREEN)

    # ----------------- SLIDE 7: Canonical WLCG Topology Model -----------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Canonical Information Model: WLCG Topology")

    add_card(s7, 0.8, 1.8, 3.6, 4.8, "1. Site Level", [
        "Unique Canonical ID (e.g. site:iisas-bratislava)",
        "Authoritative Name & Official Tier level",
        "ISO Country & Verified Geographic Coordinates",
        "Active Regional Operations Center (ROC)",
        "Participating Provider Snapshot References"
    ], ACCENT_CYAN)

    add_card(s7, 4.8, 1.8, 3.6, 4.8, "2. Service Cluster Level", [
        "CE (Computing Element): HTCondor-CE, ARC-CE",
        "SE (Storage Element): dCache, EOS, StoRM",
        "Central monitoring, squids, and site BDIIs",
        "Cluster status and operational downtime states"
    ], ACCENT_AMBER)

    add_card(s7, 8.8, 1.8, 3.6, 4.8, "3. Endpoint & VO Level", [
        "Standardized interface protocol (WLCG-SRM, WLCG-CE)",
        "FQDN, port, URL path, and security capabilities",
        "Supported Virtual Organizations (ATLAS, CMS, ALICE, LHCb)",
        "Quality assurance & health check state"
    ], ACCENT_GREEN)

    # ----------------- SLIDE 8: Infrastructure Explorer Web Interface -----------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Infrastructure Explorer: Interactive Observability UI")

    add_card(s8, 0.8, 1.8, 5.6, 4.8, "Key Frontend Capabilities", [
        "Dashboard Overview: Real-time provider health, total sites, and match ratios.",
        "Interactive World Map: D3 & React Simple Maps plotting live site coordinates.",
        "Side-by-Side Comparison: GitHub diff-style comparison of GOCDB vs BDII vs OSG.",
        "Reconciliation Inspector: Real-time step-through of evidence scoring calculation.",
        "Faceted Site Explorer: Filter by country, provider, confidence, and service type."
    ], ACCENT_CYAN)

    add_card(s8, 6.8, 1.8, 5.6, 4.8, "Built for WLCG Operations", [
        "Dark High-Tech Interface: Built with Tailwind CSS, Lucide icons, and Radix UI.",
        "Responsive & High-Performance: Sub-second search across hundreds of sites.",
        "Discrepancy Highlighting: Visual badges immediately isolate conflicting fields.",
        "Direct Provenance Tracing: View raw JSON/XML snippet from any field click.",
        "Production-grade UX ready for LHC Run 3 and HL-LHC preparation."
    ], ACCENT_PURPLE)

    # ----------------- SLIDE 9: Production Netlify Deployment -----------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Deployment: Netlify Edge & Serverless Architecture")

    add_card(s9, 0.8, 1.8, 5.6, 4.8, "Production Build & Hosting", [
        "Engineered with TanStack Start & Nitro Engine.",
        "Configured with native netlify.toml preset for automated deployment.",
        "Client static assets hosted on Netlify Global Edge CDN (dist/).",
        "Serverless functions execute SSR and API endpoints (.netlify/functions-internal).",
        "Continuous Deployment directly linked to GitHub repository."
    ], ACCENT_GREEN)

    add_card(s9, 6.8, 1.8, 5.6, 4.8, "Reliability & Scalability Benefits", [
        "Instant rollbacks and preview deployments on every Git pull request.",
        "Immutable asset caching with 1-year Cache-Control headers.",
        "Low latency worldwide for distributed WLCG computing collaborators.",
        "Zero-maintenance infrastructure scaling automatically with traffic spikes."
    ], ACCENT_CYAN)

    # ----------------- SLIDE 10: Conclusion & Roadmap -----------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "Summary & Future Roadmap for WLCG CRIC")

    add_card(s10, 0.8, 1.8, 5.6, 4.8, "POC Key Achievements", [
        "Proved automated multi-catalogue federation across GOCDB, BDII, and OSG.",
        "Achieved zero data loss via defensive valid_with_warnings classification.",
        "Delivered 100% auditable evidence scoring without black-box models.",
        "Built a modern, reactive visual explorer deployed to cloud edge.",
        "Directly aligns with CERN IT-CE-LCG mandate for reliable grid metadata."
    ], ACCENT_GREEN)

    add_card(s10, 6.8, 1.8, 5.6, 4.8, "Next Steps for Production", [
        "Scheduled Snapshot Polling: Automated 5-minute delta synchronizer.",
        "CERN IAM Integration: Role-based access control for site administrators.",
        "Historical Drift Tracking: Temporal diffs showing site changes over months.",
        "HL-LHC Scale Testing: Benchmarking against full High-Luminosity datasets."
    ], ACCENT_CYAN)

    # Save presentation
    output_path = os.path.abspath("WLCG_CRIC_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
